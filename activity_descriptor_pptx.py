from __future__ import annotations

from copy import deepcopy
from datetime import datetime
from io import BytesIO
from pathlib import Path

import math

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt


TEMPLATE_PATH = Path(__file__).resolve().parent / "templates" / "activity_descriptor.pptx"

LINKS_TEXT = (
    "Language and Literacy\n"
    "Wellness\n"
    "Creative Expression\n"
    "Mathematics\n"
    "Social Emotional"
)

FONT_PT = 12.0
MIN_FONT_PT = 8.0
LINE_HEIGHT_FACTOR = 1.1
LINE_SPACING_AFTER_PT = 3.0
AVG_CHAR_WIDTH_FACTOR = 0.40  # measured avg advance width / size for the template's font (Tw Cen MT); 0.55 was an untuned guess that shrank text well below what actually fits
CONTENT_FONT = "Times New Roman"
TIMES_CHAR_WIDTH_FACTOR = 0.48  # Times New Roman is wider than Tw Cen MT
PANEL_MIN_HEIGHT_EMU = 900000  # filters out small decorative rectangles, keeps the big background cards


def _panel_bottom_emu(shape, layout_shapes) -> int:
    """The colored background card a text box sits on (drawn on the slide layout) can be a different size than the
    text box's own xfrm — the box is allowed to be looser than the card so it doesn't clip mid-character, but that
    means fitting text purely to the box's own height can still visibly overflow the card. Find the card and use
    whichever bottom edge is tighter."""
    shape_left, shape_top, shape_width = shape.left, shape.top, shape.width
    shape_right = shape_left + shape_width
    for cand in layout_shapes:
        if None in (cand.left, cand.top, cand.width, cand.height):
            continue
        if cand.height < PANEL_MIN_HEIGHT_EMU:
            continue
        cand_right = cand.left + cand.width
        if cand.left <= shape_left and cand_right >= shape_right and cand.top <= shape_top <= cand.top + cand.height:
            return min(shape_top + shape.height, cand.top + cand.height)
    return shape_top + shape.height


def _fit_font_size_pt(shape, lines: list[str], layout_shapes, char_width_factor: float = AVG_CHAR_WIDTH_FACTOR) -> float:
    """Shrink below FONT_PT only as much as needed for all lines to fit the shape, like PowerPoint's "shrink text on overflow" — but computed and baked into the XML directly, since viewers don't reliably recompute normAutofit on their own."""
    tf = shape.text_frame
    available_width = max((shape.width - (tf.margin_left or 0) - (tf.margin_right or 0)) / 12700.0, 1.0)
    bottom_emu = _panel_bottom_emu(shape, layout_shapes)
    available_height = max((bottom_emu - shape.top - (tf.margin_top or 0) - (tf.margin_bottom or 0)) / 12700.0, 1.0)

    size = FONT_PT
    while size > MIN_FONT_PT:
        chars_per_line = max(1, int(available_width / (size * char_width_factor)))
        total_height = sum(
            max(1, math.ceil(len(line) / chars_per_line)) * size * LINE_HEIGHT_FACTOR + LINE_SPACING_AFTER_PT
            for line in lines
        )
        if total_height <= available_height:
            return size
        size -= 0.5
    return MIN_FONT_PT


_BULLET_TAGS = (
    "a:buClrTx", "a:buClr", "a:buSzTx", "a:buSzPct", "a:buSzPts",
    "a:buFontTx", "a:buFont", "a:buNone", "a:buAutoNum", "a:buChar", "a:buBlip",
)


def _remove_bullet(p_elem) -> None:
    """Strip bullet formatting and hanging indent so generated lines render as plain, flush-left text."""
    pPr = p_elem.get_or_add_pPr()
    for tag in _BULLET_TAGS:
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    pPr.set("marL", "0")
    pPr.set("indent", "0")
    insert_index = len(pPr)
    for i, child in enumerate(pPr):
        if child.tag in (qn("a:tabLst"), qn("a:defRPr"), qn("a:extLst")):
            insert_index = i
            break
    pPr.insert(insert_index, pPr.makeelement(qn("a:buNone"), {}))


def _set_text_box(shape, text: str, layout_shapes, size_pt: float | None = None, font_name: str | None = None) -> None:
    """Replace text in a shape's text frame, preserving first-run formatting."""
    tf = shape.text_frame

    # Grab formatting from the first run of the first paragraph if it exists
    first_run = None
    if tf.paragraphs and tf.paragraphs[0].runs:
        first_run = tf.paragraphs[0].runs[0]

    # Clear all existing paragraphs except the first
    for para in tf.paragraphs[1:]:
        p_elem = para._p
        p_elem.getparent().remove(p_elem)

    lines = text.split("\n")
    first_para = tf.paragraphs[0]

    # Clear runs from the first paragraph
    for run in first_para.runs:
        r_elem = run._r
        r_elem.getparent().remove(r_elem)

    # Add the first line to the existing paragraph
    run0 = first_para.add_run()
    run0.text = lines[0]
    if first_run:
        run0.font.bold = first_run.font.bold
    run0.font.size = Pt(size_pt if size_pt is not None else _fit_font_size_pt(shape, lines, layout_shapes))
    if font_name:
        run0.font.name = font_name
        # Same line spacing everywhere so no box looks tighter than another
        first_para.line_spacing = LINE_HEIGHT_FACTOR
        first_para.space_after = Pt(LINE_SPACING_AFTER_PT)

    first_p_elem = first_para._p
    _remove_bullet(first_p_elem)
    parent = first_p_elem.getparent()

    # Add subsequent lines as new paragraphs

    for line in lines[1:]:
        new_p = deepcopy(first_p_elem)
        # Remove all runs from the copied element and add fresh one
        for r in new_p.findall(qn("a:r")):
            new_p.remove(r)
        run_elem = deepcopy(first_p_elem.findall(qn("a:r"))[0]) if first_p_elem.findall(qn("a:r")) else etree.SubElement(new_p, qn("a:r"))
        # Set the text
        t_elem = run_elem.find(qn("a:t"))
        if t_elem is None:
            t_elem = etree.SubElement(run_elem, qn("a:t"))
        t_elem.text = line
        new_p.append(run_elem)
        parent.append(new_p)


def _ordinal_suffix(day: int) -> str:
    if 11 <= day % 100 <= 13:
        return "th"
    return {1: "st", 2: "nd", 3: "rd"}.get(day % 10, "th")


def _format_date_long(date_str: str) -> str:
    if not date_str:
        return ""
    try:
        date = datetime.strptime(date_str, "%Y-%m-%d")
    except ValueError:
        return date_str
    return f"{date.strftime('%b')} {date.day}{_ordinal_suffix(date.day)}, {date.year}"


def build_activity_descriptor_pptx(payload: dict) -> tuple[bytes, str]:
    raw_date = str(payload.get("date", ""))
    date_str = _format_date_long(raw_date)
    activities = payload.get("activities", [])
    skills = payload.get("skills", [])
    links = payload.get("links", [])

    content_texts = {
        "Text Box 19": "\n".join(a for a in activities if a),
        "Text Box 20": "\n".join(f"- {s}" for s in skills if s),
        "Text Box 21": "\n".join(f"• {l}" for l in links if l),
    }

    prs = Presentation(str(TEMPLATE_PATH))
    slide = prs.slides[0]
    layout_shapes = list(slide.slide_layout.shapes)

    # One shared size for all content boxes: whatever the fullest box needs.
    shared_size = FONT_PT
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name in content_texts and content_texts[shape.name]:
            shared_size = min(
                shared_size,
                _fit_font_size_pt(shape, content_texts[shape.name].split("\n"), layout_shapes, TIMES_CHAR_WIDTH_FACTOR),
            )

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        name = shape.name
        if name == "Text Box 18":
            _set_text_box(shape, date_str, layout_shapes)
        elif name in content_texts:
            if name == "Text Box 21" and not links:
                continue
            _set_text_box(shape, content_texts[name], layout_shapes, shared_size, CONTENT_FONT)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    filename = f"activity_descriptor_{raw_date or 'export'}.pptx"
    return buf.read(), filename
